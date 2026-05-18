"""
Build PowerPoint presentation summarizing the tool wear pipeline chronologically.

Sections:
  1. Title + pipeline overview
  2. Data Quality Audit
  3. Baseline (N_ST)
  4. Tuning Effect
  5. Data Augmentation
  6. Branch-Level Comparison
  7. Sequential Story (evolution)
  8. Predictions Analysis
  9. SHAP - Interpretability
  10. Conclusions
"""
from __future__ import annotations
from pathlib import Path
from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ---------------------------------------------------------------------------
# Palette (industrial: navy + steel + rust)
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
NAVY_DARK = RGBColor(0x14, 0x28, 0x46)
STEEL = RGBColor(0x5B, 0x7A, 0x9C)
RUST = RGBColor(0xD9, 0x74, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE = RGBColor(0xF8, 0xF9, 0xFB)
INK = RGBColor(0x1F, 0x29, 0x37)
MUTED = RGBColor(0x6B, 0x72, 0x80)
LIGHTGRAY = RGBColor(0xE5, 0xE7, 0xEB)
GREEN_OK = RGBColor(0x1B, 0x7F, 0x5A)
RED_BAD = RGBColor(0xB0, 0x32, 0x4A)


BASE = Path(__file__).parent.parent
FIG = BASE / 'outputs' / 'figures'
OUT_PPTX = BASE / 'outputs' / 'reports' / 'tool_wear_pipeline_Manuel_Pusma.pptx'
OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def hexc(r, g, b):
    return RGBColor(r, g, b)


def img_aspect(path: Path) -> float:
    with Image.open(path) as im:
        w, h = im.size
    return w / h


def fit_image(slide, path: Path, left, top, max_w, max_h):
    """Add an image fitting within (max_w, max_h) preserving aspect."""
    ar = img_aspect(path)
    box_ar = max_w / max_h
    if ar > box_ar:
        w = max_w
        h = max_w / ar
    else:
        h = max_h
        w = max_h * ar
    cx = left + (max_w - w) / 2
    cy = top + (max_h - h) / 2
    slide.shapes.add_picture(str(path), Inches(cx), Inches(cy),
                             width=Inches(w), height=Inches(h))


def add_rect(slide, left, top, w, h, fill_rgb, line_rgb=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(left), Inches(top),
                                  Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_rgb
    if line_rgb is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_rgb
    shp.shadow.inherit = False
    return shp


def add_text(slide, text, left, top, w, h,
             font='Calibri', size=14, color=INK, bold=False,
             italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split('\n') if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.italic = italic
    return tb


def add_bullets(slide, items, left, top, w, h,
                font='Calibri', size=14, color=INK, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # Bullet via leading character + space; cleaner than xml bullets here
        r = p.add_run()
        r.text = f'•  {item}'
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def add_full_bg(slide, color):
    """Fill entire slide background."""
    add_rect(slide, 0, 0, SW, SH, color)


def slide_title(slide, title, sub=None, color=INK,
                sub_color=MUTED):
    add_text(slide, title, 0.5, 0.35, SW - 1.0, 0.7,
             font='Calibri', size=28, color=color, bold=True)
    if sub:
        add_text(slide, sub, 0.5, 1.05, SW - 1.0, 0.45,
                 font='Calibri', size=15, color=sub_color, italic=True)


# ---------------------------------------------------------------------------
# Presentation setup
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SW = 13.333
SH = 7.5

blank = prs.slide_layouts[6]


# =============================================================================
# SLIDE 1 — TITLE
# =============================================================================
s = prs.slides.add_slide(blank)
add_full_bg(s, NAVY)

# Decorative thin rust band at the bottom for branding
add_rect(s, 0, SH - 0.35, SW, 0.35, RUST)

add_text(s, 'TOOL WEAR PREDICTION', 0.8, 2.0, SW - 1.6, 1.0,
         font='Calibri', size=54, color=WHITE, bold=True)
add_text(s, 'A layered methodology study on the PHM tool wear dataset',
         0.8, 3.1, SW - 1.6, 0.6, font='Calibri', size=22,
         color=hexc(0xCB, 0xD5, 0xE1), italic=True)

# Tagline / metadata block
add_rect(s, 0.8, 4.5, 4.5, 1.5, NAVY_DARK)
add_text(s, 'Dataset', 1.0, 4.65, 4.0, 0.4, size=11,
         color=hexc(0xCB, 0xD5, 0xE1), bold=True)
add_text(s, '10 experiments  ·  1 tool (T01)  ·  Target: VB_µm',
         1.0, 4.95, 4.0, 0.5, size=15, color=WHITE, bold=False)
add_text(s, 'Validation', 1.0, 5.45, 4.0, 0.3, size=11,
         color=hexc(0xCB, 0xD5, 0xE1), bold=True)
add_text(s, 'LOEO-CV  ·  10 folds  ·  honest', 1.0, 5.7, 4.0, 0.4,
         size=15, color=WHITE)

add_rect(s, 5.8, 4.5, 4.5, 1.5, NAVY_DARK)
add_text(s, 'Best model', 6.0, 4.65, 4.0, 0.4, size=11,
         color=hexc(0xCB, 0xD5, 0xE1), bold=True)
add_text(s, 'ElasticNet  →  MAE = 18.79 µm  ·  R² = 0.82', 6.0, 4.95, 4.0, 0.5,
         size=15, color=WHITE, bold=True)
add_text(s, 'Winning branch', 6.0, 5.45, 4.0, 0.3, size=11,
         color=hexc(0xCB, 0xD5, 0xE1), bold=True)
add_text(s, 'SOLO_A_N_CT_random (axial-only, tuned, no augmentation)',
         6.0, 5.7, 4.0, 0.4, size=13, color=WHITE)


# =============================================================================
# SLIDE 2 — PIPELINE OVERVIEW (flow diagram)
# =============================================================================
s = prs.slides.add_slide(blank)
add_full_bg(s, WHITE)
slide_title(s, 'The pipeline at a glance',
            sub='From raw data D, three feature subsets (FUSION / SOLO_A / SOLO_R) × '
                'two data branches (Normal / Augmented) × three tuning options = '
                '36 branches under one LOEO-CV evaluator, ranking, then SHAP.')
fit_image(s, FIG / 'layered_pipeline' / '00_layered_flow_diagram_no_holdout.png',
          0.5, 1.55, SW - 1.0, SH - 2.0)


# =============================================================================
# SLIDE 3 — SECTION DIVIDER: Data Quality
# =============================================================================
def section_divider(num, title, blurb):
    s = prs.slides.add_slide(blank)
    add_full_bg(s, NAVY)
    add_rect(s, 0, SH - 0.35, SW, 0.35, RUST)
    add_text(s, f'STEP {num}', 0.8, 2.4, 6.0, 0.6, size=18,
             color=RUST, bold=True)
    add_text(s, title, 0.8, 2.9, SW - 1.6, 1.2, size=48,
             color=WHITE, bold=True)
    add_text(s, blurb, 0.8, 4.4, SW - 1.6, 1.5, size=18,
             color=hexc(0xCB, 0xD5, 0xE1), italic=True)
    return s


section_divider(
    1, 'Data Quality Audit',
    'Before training any model: what do we actually have on disk, and is '
    'the target variable usable? Ten experiments, one tool, contact-mounted '
    'measurements.')


# =============================================================================
# SLIDE 4 — Data quality: raw count + missing segments
# =============================================================================
s = prs.slides.add_slide(blank)
add_full_bg(s, WHITE)
slide_title(s, 'What is on disk',
            sub='Per-experiment file inventory and missing-segment audit.')

# Two panels
fit_image(s, FIG / 'data_quality' / 'raw_file_count_by_experiment.png',
          0.4, 1.55, 6.3, 5.5)
fit_image(s, FIG / 'data_quality' / 'missing_segments_by_experiment.png',
          6.7, 1.55, 6.3, 5.5)

# Captions
add_text(s, 'Files per experiment', 0.4, SH - 0.45, 6.3, 0.3,
         size=11, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
add_text(s, 'Missing segments by experiment', 6.7, SH - 0.45, 6.3, 0.3,
         size=11, color=MUTED, italic=True, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 5 — Data quality: target distribution + order
# =============================================================================
s = prs.slides.add_slide(blank)
add_full_bg(s, WHITE)
slide_title(s, 'The target variable VB_µm',
            sub='Range, distribution and trajectory across the 10 experiments. '
                'No replicates per experiment — each VB is a single measurement.')

fit_image(s, FIG / 'data_quality' / 'vb_distribution.png',
          0.4, 1.55, 6.3, 5.3)
fit_image(s, FIG / 'data_quality' / 'vb_vs_experiment_order.png',
          6.7, 1.55, 6.3, 5.3)

add_text(s, 'VB_µm distribution (n=10)', 0.4, SH - 0.4, 6.3, 0.3,
         size=11, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
add_text(s, 'VB_µm versus experiment order', 6.7, SH - 0.4, 6.3, 0.3,
         size=11, color=MUTED, italic=True, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 6 — SECTION DIVIDER: Baseline
# =============================================================================
section_divider(
    2, 'Baseline — N_ST',
    'Normal data, no tuning, no augmentation. Eight regressors trained '
    'with default hyper-parameters. This is the bar everything else must beat.')


# =============================================================================
# SLIDE 7 — Baseline narrative (numbers + branch_performance preview)
# =============================================================================
s = prs.slides.add_slide(blank)
add_full_bg(s, WHITE)
slide_title(s, 'Baseline performance (FUSION_N_ST)',
            sub='ElasticNet emerges as the strongest baseline. Tree-based '
                'and kernel models underperform with only 8 training rows.')

# Left: big stat callout
add_rect(s, 0.5, 1.7, 5.2, 4.8, OFFWHITE)
add_text(s, 'BASELINE  (FUSION_N_ST)', 0.7, 1.85, 5.0, 0.4, size=14,
         color=STEEL, bold=True)
add_text(s, 'MAE  26.96 µm', 0.7, 2.3, 5.0, 0.9, size=42,
         color=NAVY, bold=True)
add_text(s, 'RMSE  37.88 µm', 0.7, 3.2, 5.0, 0.55, size=20,
         color=INK)
add_text(s, 'R²  0.602', 0.7, 3.75, 5.0, 0.55, size=20, color=INK)
add_text(s, 'MAPE  17.0 %', 0.7, 4.3, 5.0, 0.55, size=20, color=INK)
add_text(s, 'winning model: ElasticNet', 0.7, 5.0, 5.0, 0.4,
         size=13, color=MUTED, italic=True)
add_text(s, 'config: fusion (203 feats: A+R+agg), no tuning, no augmentation',
         0.7, 5.35, 5.0, 0.4, size=12, color=MUTED, italic=True)
add_text(s, 'training rows: 8  ·  test rows: 1 per fold  ·  10 folds',
         0.7, 5.7, 5.0, 0.4, size=13, color=MUTED, italic=True)

# Right: branch_performance_MAE image
fit_image(s, FIG / 'layered_pipeline' / '09_branch_performance_MAE.png',
          6.0, 1.6, 6.9, 5.5)
add_text(s, 'MAE by branch — N_ST highlighted in context',
         6.0, SH - 0.4, 6.9, 0.3, size=11, color=MUTED,
         italic=True, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 8 — SECTION DIVIDER: Tuning
# =============================================================================
section_divider(
    3, 'Tuning Effect — Does It Help?',
    'RandomizedSearchCV vs GridSearchCV on top of normal data. '
    'Search is performed once over the train pool with internal GroupKFold; '
    'best params are then reused inside LOEO folds.')


# =============================================================================
# SLIDE 9 — Tuning effect MAE + RMSE
# =============================================================================
s = prs.slides.add_slide(blank)
add_full_bg(s, WHITE)
slide_title(s, 'Tuning effect by data branch',
            sub='Random vs Grid vs no-tuning. With n=8 train rows, search spaces '
                'are intentionally small — tuning can only do so much.')

fit_image(s, FIG / 'layered_pipeline' / '09_tuning_effect_MAE.png',
          0.4, 1.55, 6.3, 5.4)
fit_image(s, FIG / 'layered_pipeline' / '09_tuning_effect_RMSE.png',
          6.7, 1.55, 6.3, 5.4)

add_text(s, 'MAE', 0.4, SH - 0.4, 6.3, 0.3, size=11, color=MUTED,
         italic=True, align=PP_ALIGN.CENTER)
add_text(s, 'RMSE', 6.7, SH - 0.4, 6.3, 0.3, size=11, color=MUTED,
         italic=True, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 10 — Random vs Grid
# =============================================================================
s = prs.slides.add_slide(blank)
add_full_bg(s, WHITE)
slide_title(s, 'Random vs Grid search',
            sub='With such small search spaces the two methods converge on the '
                'same hyper-parameters — Random and Grid produce identical '
                'metrics for most models. Grid did not find anything Random missed.')

fit_image(s, FIG / 'layered_pipeline' / '09_random_vs_grid_MAE.png',
          1.5, 1.55, 10.3, 5.4)
add_text(s, 'MAE — Random vs Grid across models',
         1.5, SH - 0.4, 10.3, 0.3, size=11, color=MUTED,
         italic=True, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 11 — SECTION DIVIDER: Augmentation
# =============================================================================
section_divider(
    4, 'Data Augmentation',
    'Three perturbation strategies applied to training rows only: '
    'feature_noise (Gaussian), feature_scaling (multiplicative), '
    'grouped_scaling (per feature-group). Real test rows are never touched.')


# =============================================================================
# SLIDE 12 — Augmentation effect MAE + RMSE
# =============================================================================
s = prs.slides.add_slide(blank)
add_full_bg(s, WHITE)
slide_title(s, 'Augmentation effect by strategy',
            sub='Augmentation alone (A_ST) is essentially neutral on MAE. '
                'feature_scaling visibly reduces RMSE — outliers shrink.')

fit_image(s, FIG / 'layered_pipeline' / '09_augmentation_effect_MAE.png',
          0.4, 1.55, 6.3, 5.4)
fit_image(s, FIG / 'layered_pipeline' / '09_augmentation_effect_RMSE.png',
          6.7, 1.55, 6.3, 5.4)
add_text(s, 'MAE', 0.4, SH - 0.4, 6.3, 0.3, size=11, color=MUTED,
         italic=True, align=PP_ALIGN.CENTER)
add_text(s, 'RMSE', 6.7, SH - 0.4, 6.3, 0.3, size=11, color=MUTED,
         italic=True, align=PP_ALIGN.CENTER)


# =============================================================================
# NEW — SECTION DIVIDER: Feature Subset Branching (FUSION/SOLO_A/SOLO_R)
# =============================================================================
section_divider(
    5, 'Feature-Subset Branching',
    'Splitting the tree at the very top by signal type: '
    'FUSION (all 203 features), SOLO_A (axial only, ~101), '
    'SOLO_R (rotational only, ~99). Motivation: EDA showed '
    'mean |corr(A, R)| ≈ 0.70 — A and R encode largely the same wear signal.')


# =============================================================================
# NEW — EDA redundancy (justification for SOLO branches)
# =============================================================================
s = prs.slides.add_slide(blank)
add_full_bg(s, WHITE)
slide_title(s, 'Why split into SOLO_A and SOLO_R?',
            sub='Axial (A) and rotational (R) features share most of the '
                'wear signal — redundancy is high enough that the model can '
                'learn from either one alone.')

fit_image(s, FIG / 'eda_fusion' / 'redundancy_heatmap.png',
          0.4, 1.55, 6.3, 5.4)
fit_image(s, FIG / 'eda_fusion' / 'top_features_per_direction.png',
          6.7, 1.55, 6.3, 5.4)
add_text(s, 'A vs R redundancy heatmap', 0.4, SH - 0.4, 6.3, 0.3,
         size=11, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
add_text(s, 'Top correlated features per direction',
         6.7, SH - 0.4, 6.3, 0.3, size=11, color=MUTED, italic=True,
         align=PP_ALIGN.CENTER)


# =============================================================================
# NEW — Signal branch experiment outcome (3 best models)
# =============================================================================
s = prs.slides.add_slide(blank)
add_full_bg(s, WHITE)
slide_title(s, 'SOLO_A beats FUSION by ~30 % MAE',
            sub='Isolated baseline comparison (no tuning, no augmentation) '
                'across all baseline regressors. SOLO_A wins consistently — '
                'removing the redundant R features cleans the signal.')

fit_image(s, FIG / 'signal_branch' / 'best_model_per_branch.png',
          0.4, 1.55, 8.3, 5.4)

# Right: stat callout per branch
add_rect(s, 9.0, 1.7, 4.0, 5.0, OFFWHITE)
add_text(s, 'PER-BRANCH BEST', 9.2, 1.85, 3.6, 0.4, size=12,
         color=STEEL, bold=True)
add_text(s, 'FUSION', 9.2, 2.3, 3.6, 0.35, size=14, color=NAVY, bold=True)
add_text(s, 'MAE 30.2 µm  ·  203 feats', 9.2, 2.65, 3.6, 0.35,
         size=12, color=INK)
add_text(s, 'SOLO_A', 9.2, 3.2, 3.6, 0.35, size=14,
         color=hexc(0xA0, 0x52, 0x1E), bold=True)
add_text(s, 'MAE 21.6 µm  ·  101 feats', 9.2, 3.55, 3.6, 0.35,
         size=12, color=INK)
add_text(s, 'SOLO_R', 9.2, 4.1, 3.6, 0.35, size=14,
         color=hexc(0x1B, 0x7F, 0x5A), bold=True)
add_text(s, 'MAE 30.9 µm  ·  99 feats', 9.2, 4.45, 3.6, 0.35,
         size=12, color=INK)
add_text(s, 'Verdict', 9.2, 5.1, 3.6, 0.35, size=12, color=MUTED,
         italic=True, bold=True)
add_text(s, 'Drop R features.\nWear lives in axial.',
         9.2, 5.45, 3.6, 0.9, size=13, color=INK, bold=True)


# =============================================================================
# SLIDE 13 — SECTION DIVIDER: Branch comparison
# =============================================================================
section_divider(
    6, 'Branch-Level Comparison',
    'All thirty-six branches side by side (3 feature subsets × 12 stages): '
    'per-branch best, deltas against baseline, and full branch × model heatmaps.')


# =============================================================================
# SLIDE 14 — Delta MAE + RMSE vs baseline
# =============================================================================
s = prs.slides.add_slide(blank)
add_full_bg(s, WHITE)
slide_title(s, 'Δ versus baseline FUSION_N_ST',
            sub='Negative = improves. SOLO_A branches dominate the green band '
                '(-5 to -8 µm MAE). FUSION_A_CT_*_feature_noise spikes (Lasso '
                'collapses); SOLO_R branches widen the error.')

fit_image(s, FIG / 'layered_pipeline' / '09_delta_MAE_vs_baseline_FUSION_N_ST.png',
          0.4, 1.55, 6.3, 5.4)
fit_image(s, FIG / 'layered_pipeline' / '09_delta_RMSE_vs_baseline_FUSION_N_ST.png',
          6.7, 1.55, 6.3, 5.4)
add_text(s, 'Δ MAE', 0.4, SH - 0.4, 6.3, 0.3, size=11, color=MUTED,
         italic=True, align=PP_ALIGN.CENTER)
add_text(s, 'Δ RMSE', 6.7, SH - 0.4, 6.3, 0.3, size=11, color=MUTED,
         italic=True, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 15 — Heatmap MAE
# =============================================================================
s = prs.slides.add_slide(blank)
add_full_bg(s, WHITE)
slide_title(s, 'Branch × model heatmap — MAE',
            sub='Linear-regularized models (Ridge, Lasso, ElasticNet) dominate. '
                'SOLO_A rows score lowest; SOLO_R rows score worst. '
                'Tree ensembles and MLP suffer the small-sample curse.')

fit_image(s, FIG / 'layered_pipeline' / '09_heatmap_model_vs_branch_MAE.png',
          0.5, 1.55, SW - 1.0, 5.4)


# =============================================================================
# SLIDE 16 — Heatmap R²
# =============================================================================
s = prs.slides.add_slide(blank)
add_full_bg(s, WHITE)
slide_title(s, 'Branch × model heatmap — R²',
            sub='Same shape as MAE. SOLO_A_N_CT_random with ElasticNet '
                'reaches R² ≈ 0.82; tree ensembles and MLP often go negative.')

fit_image(s, FIG / 'layered_pipeline' / '09_heatmap_model_vs_branch_R2.png',
          0.5, 1.55, SW - 1.0, 5.4)


# =============================================================================
# SLIDE 17 — Best per branch
# =============================================================================
s = prs.slides.add_slide(blank)
add_full_bg(s, WHITE)
slide_title(s, 'Best model per branch',
            sub='ElasticNet wins the vast majority of the 36 branches. '
                'SOLO_A_N_CT_random / _grid sit at the top (MAE 18.79 µm); '
                'a few FUSION_A_CT_*_feature_noise branches collapse onto Lasso.')

fit_image(s, FIG / 'layered_pipeline' / '09_best_model_per_branch_MAE.png',
          0.5, 1.55, SW - 1.0, 5.4)


# =============================================================================
# SLIDE 18 — Branch performance MAE + RMSE
# =============================================================================
s = prs.slides.add_slide(blank)
add_full_bg(s, WHITE)
slide_title(s, 'Branch performance — MAE & RMSE',
            sub='Lower is better. The 36 branches separate into three groups '
                'by feature subset: SOLO_A (best, blue band), FUSION (middle), '
                'SOLO_R (worst). SOLO_A_N_CT_random tops at 18.79 µm.')

fit_image(s, FIG / 'layered_pipeline' / '09_branch_performance_MAE.png',
          0.4, 1.55, 6.3, 5.4)
fit_image(s, FIG / 'layered_pipeline' / '09_branch_performance_RMSE.png',
          6.7, 1.55, 6.3, 5.4)
add_text(s, 'MAE', 0.4, SH - 0.4, 6.3, 0.3, size=11, color=MUTED,
         italic=True, align=PP_ALIGN.CENTER)
add_text(s, 'RMSE', 6.7, SH - 0.4, 6.3, 0.3, size=11, color=MUTED,
         italic=True, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 19 — Branch performance R² + MAPE
# =============================================================================
s = prs.slides.add_slide(blank)
add_full_bg(s, WHITE)
slide_title(s, 'Branch performance — R² & MAPE',
            sub='R² and MAPE confirm the MAE ranking. SOLO_A_N_CT_* reaches '
                'R² ≈ 0.82, MAPE ≈ 14 %. SOLO_R consistently collapses to '
                'R² near zero or negative.')

fit_image(s, FIG / 'layered_pipeline' / '09_branch_performance_R2.png',
          0.4, 1.55, 6.3, 5.4)
fit_image(s, FIG / 'layered_pipeline' / '09_branch_performance_MAPE.png',
          6.7, 1.55, 6.3, 5.4)
add_text(s, 'R²  (higher is better)', 0.4, SH - 0.4, 6.3, 0.3,
         size=11, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
add_text(s, 'MAPE  (lower is better)', 6.7, SH - 0.4, 6.3, 0.3,
         size=11, color=MUTED, italic=True, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 20 — SECTION DIVIDER: Sequential Story
# =============================================================================
section_divider(
    6, 'Sequential Story',
    'Following a practitioner\'s mental walk: baseline → tuning → augmentation '
    '→ both. Did each layer of complexity pay off?')


# =============================================================================
# SLIDE 21 — Sequential dashboard MAE (key viz)
# =============================================================================
s = prs.slides.add_slide(blank)
add_full_bg(s, WHITE)
slide_title(s, 'Sequential dashboard — MAE',
            sub='Four panels showing the four moves a practitioner would try. '
                'On MAE, the gains across stages are within noise (n=10).')

fit_image(s, FIG / 'layered_pipeline' / '09_sequential_comparison_dashboard_MAE.png',
          0.4, 1.55, SW - 0.8, 5.5)


# =============================================================================
# SLIDE 22 — Sequential dashboard RMSE + R²
# =============================================================================
s = prs.slides.add_slide(blank)
add_full_bg(s, WHITE)
slide_title(s, 'Sequential dashboard — RMSE & R²',
            sub='On RMSE the picture changes: tuning + augmentation does cut '
                'tail error, and R² improves visibly when both are combined.')

fit_image(s, FIG / 'layered_pipeline' / '09_sequential_comparison_dashboard_RMSE.png',
          0.4, 1.55, 6.3, 5.4)
fit_image(s, FIG / 'layered_pipeline' / '09_sequential_comparison_dashboard_R2.png',
          6.7, 1.55, 6.3, 5.4)
add_text(s, 'RMSE', 0.4, SH - 0.4, 6.3, 0.3, size=11, color=MUTED,
         italic=True, align=PP_ALIGN.CENTER)
add_text(s, 'R²', 6.7, SH - 0.4, 6.3, 0.3, size=11, color=MUTED,
         italic=True, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 23 — Model evolution MAE + RMSE
# =============================================================================
s = prs.slides.add_slide(blank)
add_full_bg(s, WHITE)
slide_title(s, 'Model evolution — MAE & RMSE',
            sub='Six progressive stages on the same model. Δ between '
                'consecutive stages annotated (green=better, red=worse).')

fit_image(s, FIG / 'layered_pipeline' / '09_model_evolution_MAE_LOEO.png',
          0.4, 1.55, 6.3, 5.4)
fit_image(s, FIG / 'layered_pipeline' / '09_model_evolution_RMSE_LOEO.png',
          6.7, 1.55, 6.3, 5.4)
add_text(s, 'MAE', 0.4, SH - 0.4, 6.3, 0.3, size=11, color=MUTED,
         italic=True, align=PP_ALIGN.CENTER)
add_text(s, 'RMSE', 6.7, SH - 0.4, 6.3, 0.3, size=11, color=MUTED,
         italic=True, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 24 — Model evolution R² + MAPE
# =============================================================================
s = prs.slides.add_slide(blank)
add_full_bg(s, WHITE)
slide_title(s, 'Model evolution — R² & MAPE',
            sub='R² improves clearly once tuning is added on augmented data. '
                'MAPE stays remarkably stable (~17 %), suggesting that the '
                'percentage error floor is data-driven, not model-driven.')

fit_image(s, FIG / 'layered_pipeline' / '09_model_evolution_R2_LOEO.png',
          0.4, 1.55, 6.3, 5.4)
fit_image(s, FIG / 'layered_pipeline' / '09_model_evolution_MAPE_LOEO.png',
          6.7, 1.55, 6.3, 5.4)
add_text(s, 'R²', 0.4, SH - 0.4, 6.3, 0.3, size=11, color=MUTED,
         italic=True, align=PP_ALIGN.CENTER)
add_text(s, 'MAPE  (%)', 6.7, SH - 0.4, 6.3, 0.3, size=11, color=MUTED,
         italic=True, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 25 — SECTION DIVIDER: Predictions
# =============================================================================
section_divider(
    7, 'Predictions Analysis',
    'How the predicted VB values line up against the truth, fold by fold. '
    'Multi-configuration overlays first, then a deeper look at the best '
    'global model.')


# =============================================================================
# SLIDE 26 — Multi-config: actual vs predicted
# =============================================================================
s = prs.slides.add_slide(blank)
add_full_bg(s, WHITE)
slide_title(s, 'Actual vs predicted — five configurations overlaid',
            sub='Five colors trace baseline, tuned, augmented, augmented+tuned, '
                'and best global. Predictions cluster in the same diagonal — '
                'the gain from each stage is small.')

fit_image(s, FIG / 'layered_pipeline' / '09_actual_vs_predicted_multi_LOEO.png',
          0.5, 1.55, SW - 1.0, 5.5)


# =============================================================================
# SLIDE 27 — Multi-config: residuals by experiment
# =============================================================================
s = prs.slides.add_slide(blank)
add_full_bg(s, WHITE)
slide_title(s, 'Residuals by experiment — five configurations',
            sub='Experiment-by-experiment view: the largest residuals occur on '
                'a small number of experiments (likely outliers in VB range) — '
                'all configurations miss the same ones.')

fit_image(s, FIG / 'layered_pipeline' / '09_residuals_by_experiment_multi_LOEO.png',
          0.5, 1.55, SW - 1.0, 5.5)


# =============================================================================
# SLIDE 28 — Best global: actual vs predicted + residuals scatter
# =============================================================================
s = prs.slides.add_slide(blank)
add_full_bg(s, WHITE)
slide_title(s, 'Best global model — ElasticNet on A_ST_feature_noise',
            sub='Predicted vs actual diagonal and residuals scatter for the '
                'winning configuration in LOEO.')

fit_image(s, FIG / 'layered_pipeline' / '09_actual_vs_predicted_best_global_LOEO.png',
          0.4, 1.55, 6.3, 5.4)
fit_image(s, FIG / 'layered_pipeline' / '09_residuals_best_global_LOEO.png',
          6.7, 1.55, 6.3, 5.4)
add_text(s, 'Actual vs predicted', 0.4, SH - 0.4, 6.3, 0.3, size=11,
         color=MUTED, italic=True, align=PP_ALIGN.CENTER)
add_text(s, 'Residuals vs predicted', 6.7, SH - 0.4, 6.3, 0.3, size=11,
         color=MUTED, italic=True, align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 29 — Best global: residuals by experiment
# =============================================================================
s = prs.slides.add_slide(blank)
add_full_bg(s, WHITE)
slide_title(s, 'Best global — residuals per experiment',
            sub='Each bar is one held-out experiment. The error budget is '
                'concentrated in 2–3 experiments at the extremes of the VB range.')

fit_image(s, FIG / 'layered_pipeline' / '09_residuals_by_experiment_best_global_LOEO.png',
          0.5, 1.55, SW - 1.0, 5.5)


# =============================================================================
# SLIDE 30 — SECTION DIVIDER: SHAP
# =============================================================================
section_divider(
    8, 'SHAP — Interpretability (Final Step)',
    'SHAP runs strictly after the metrics have chosen the best models. '
    'It does not train, does not tune, does not pick the winner — it only '
    'explains why the winning models predict what they predict.')


# =============================================================================
# SLIDE 31 — SHAP baseline N_ST ElasticNet
# =============================================================================
s = prs.slides.add_slide(blank)
add_full_bg(s, WHITE)
slide_title(s, 'SHAP — ElasticNet on BEST GLOBAL (SOLO_A_N_CT_random)',
            sub='Mean |SHAP| feature ranking (bar) and signed per-experiment '
                'contributions (beeswarm) for the winning configuration '
                '(MAE = 18.79 µm).')

fit_image(s, FIG / 'shap' / '10_shap_bar_elasticnet_solo_a_n_ct_random.png',
          0.4, 1.55, 6.3, 5.4)
fit_image(s, FIG / 'shap' / '10_shap_summary_elasticnet_solo_a_n_ct_random.png',
          6.7, 1.55, 6.3, 5.4)
add_text(s, 'Mean |SHAP| importance', 0.4, SH - 0.4, 6.3, 0.3,
         size=11, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
add_text(s, 'Beeswarm — signed contributions per experiment',
         6.7, SH - 0.4, 6.3, 0.3, size=11, color=MUTED, italic=True,
         align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 32 — SHAP best (A_ST feature_noise)
# =============================================================================
s = prs.slides.add_slide(blank)
add_full_bg(s, WHITE)
slide_title(s, 'SHAP — ElasticNet on twin branch (SOLO_A_N_CT_grid)',
            sub='Same model family, identical hyperparameters reached via Grid '
                'search instead of Random. Top features are virtually identical '
                'to the Random twin — ranking is robust to the search strategy.')

fit_image(s, FIG / 'shap' / '10_shap_bar_elasticnet_solo_a_n_ct_grid.png',
          0.4, 1.55, 6.3, 5.4)
fit_image(s, FIG / 'shap' / '10_shap_summary_elasticnet_solo_a_n_ct_grid.png',
          6.7, 1.55, 6.3, 5.4)
add_text(s, 'Mean |SHAP| importance', 0.4, SH - 0.4, 6.3, 0.3,
         size=11, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
add_text(s, 'Beeswarm — signed contributions per experiment',
         6.7, SH - 0.4, 6.3, 0.3, size=11, color=MUTED, italic=True,
         align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 33 — SHAP nonlinear (XGBoost on N_CT_random)
# =============================================================================
s = prs.slides.add_slide(blank)
add_full_bg(s, WHITE)
slide_title(s, 'SHAP — XGBoost on SOLO_R_A_CT_random_feature_noise (best nonlinear)',
            sub='A nonlinear sanity check, on the rotational-only branch. '
                'XGBoost surfaces partially different features and interactions, '
                'but the dominant signal overlaps with the linear models — '
                'wear cues are similar across model families.')

fit_image(s, FIG / 'shap' / '10_shap_bar_xgboost_solo_r_a_ct_random_feature_noise.png',
          0.4, 1.55, 6.3, 5.4)
fit_image(s, FIG / 'shap' / '10_shap_summary_xgboost_solo_r_a_ct_random_feature_noise.png',
          6.7, 1.55, 6.3, 5.4)
add_text(s, 'Mean |SHAP| importance', 0.4, SH - 0.4, 6.3, 0.3,
         size=11, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
add_text(s, 'Beeswarm — signed contributions per experiment',
         6.7, SH - 0.4, 6.3, 0.3, size=11, color=MUTED, italic=True,
         align=PP_ALIGN.CENTER)


# =============================================================================
# SLIDE 34 — Conclusions
# =============================================================================
s = prs.slides.add_slide(blank)
add_full_bg(s, NAVY)
add_rect(s, 0, SH - 0.35, SW, 0.35, RUST)

add_text(s, 'TAKEAWAYS', 0.8, 0.5, 6.0, 0.5, size=14,
         color=RUST, bold=True)
add_text(s, 'What this study actually shows', 0.8, 0.9, SW - 1.6, 0.8,
         size=34, color=WHITE, bold=True)

items = [
    ('1.  SOLO_A wins — drop the rotational signal',
     'Axial-only features cut baseline MAE by ~30 % (26.96 → 18.79 µm). A '
     'and R features encode the same wear signal (|corr| ≈ 0.70); fusing '
     'them just adds noise.'),
    ('2.  Linear-regularized models dominate',
     'ElasticNet wins almost all of the 36 branches. Tree ensembles and '
     'kernels are starved by n=8 training rows.'),
    ('3.  Tuning helps inside SOLO_A',
     'Random and Grid converge on identical hyper-parameters and deliver '
     '~2.8 µm MAE gain over SOLO_A_N_ST. On FUSION the effect is noisier.'),
    ('4.  Augmentation is roughly neutral',
     'feature_noise can hurt under heavy tuning (Lasso collapse). '
     'feature_scaling slightly improves RMSE. No strategy reshapes the ranking.'),
    ('5.  SHAP is the conclusion, not the engine',
     'It runs last, on real rows only, to describe — not decide — which '
     'features drive the predictions.'),
    ('6.  Honest about scale',
     'n = 10, one tool. Any |Δ| < 5 µm should be read as noise.'),
]

ty = 1.85
for hdr, body in items:
    add_text(s, hdr, 0.8, ty, SW - 1.6, 0.4, size=16,
             color=RUST, bold=True)
    add_text(s, body, 0.8, ty + 0.35, SW - 1.6, 0.5, size=13,
             color=hexc(0xE5, 0xE7, 0xEB))
    ty += 0.82


# =============================================================================
# Save
# =============================================================================
prs.save(str(OUT_PPTX))
print(f'OK: wrote {OUT_PPTX}  ({len(prs.slides)} slides)')
